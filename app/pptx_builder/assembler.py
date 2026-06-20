import io
from typing import Optional

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# Agence VU brand colors
COLOR_DARK_BLUE = RGBColor(0x1A, 0x2E, 0x4A)   # Dark navy — headers
COLOR_ACCENT = RGBColor(0x00, 0x8B, 0xD2)       # VU blue — accents
COLOR_LIGHT_GREY = RGBColor(0xF4, 0xF6, 0xF8)   # Background panels
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT = RGBColor(0x2C, 0x3E, 0x50)         # Body text
COLOR_GOOD = RGBColor(0x27, 0xAE, 0x60)         # Green — bon
COLOR_MEDIUM = RGBColor(0xF3, 0x9C, 0x12)       # Orange — moyen
COLOR_BAD = RGBColor(0xE7, 0x4C, 0x3C)          # Red — faible

# Slide dimensions — 16:9 widescreen
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)


class PPTXAssembler:
    """Assembles PowerPoint presentations from generated slide content."""

    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize assembler.

        Args:
            template_path: Optional path to a .pptx template file.
                           If None, creates a blank presentation.
        """
        self.template_path = template_path

    def _create_presentation(self) -> Presentation:
        """Create a new Presentation, from template if provided."""
        if self.template_path:
            try:
                prs = Presentation(self.template_path)
                return prs
            except Exception:
                pass
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        return prs

    def _add_cover_slide(self, prs: Presentation, pharmacy_name: str) -> None:
        """Add a cover/title slide for the Performance Globale section."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Background rectangle
        bg = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0),
            SLIDE_WIDTH, SLIDE_HEIGHT,
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_DARK_BLUE
        bg.line.fill.background()

        # Accent bar
        bar = slide.shapes.add_shape(
            1,
            Inches(0), Inches(5.8),
            Inches(4), Inches(0.15),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_ACCENT
        bar.line.fill.background()

        # Section label
        section_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11), Inches(0.8)
        )
        tf = section_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "PARTIE 2 — PERFORMANCE GLOBALE"
        run.font.size = Pt(16)
        run.font.color.rgb = COLOR_ACCENT
        run.font.bold = True
        run.font.name = "Calibri"

        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.1), Inches(11), Inches(1.4)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"Audit 360° — {pharmacy_name}" if pharmacy_name else "Audit 360°"
        run.font.size = Pt(36)
        run.font.color.rgb = COLOR_WHITE
        run.font.bold = True
        run.font.name = "Calibri"

        # Subtitle
        sub_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.6), Inches(11), Inches(0.6)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Analyse des indicateurs de performance — Agence VU"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)
        run.font.name = "Calibri"

    # ── Helpers graphiques ────────────────────────────────────────────────────

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> RGBColor:
        h = hex_color.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _add_chart_to_slide(self, slide, chart_spec: dict,
                             x: float, y: float, w: float, h: float) -> bool:
        """
        Ajoute un graphique python-pptx sur le slide selon la spec chart_builder.
        Retourne True si réussi.
        """
        try:
            chart_type_map = {
                "bar":    XL_CHART_TYPE.BAR_CLUSTERED,
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "line":   XL_CHART_TYPE.LINE,
                "pie":    XL_CHART_TYPE.PIE,
            }
            xl_type = chart_type_map.get(chart_spec.get("type", "column"),
                                          XL_CHART_TYPE.COLUMN_CLUSTERED)

            cd = ChartData()
            cd.categories = chart_spec.get("categories", [])
            for series in chart_spec.get("series", []):
                cd.add_series(series["name"], series.get("values", []))

            chart_obj = slide.shapes.add_chart(
                xl_type,
                Inches(x), Inches(y), Inches(w), Inches(h),
                cd,
            )
            chart = chart_obj.chart

            # Titre du graphique
            title_str = chart_spec.get("title", "")
            if title_str:
                chart.has_title = True
                chart.chart_title.text_frame.text = title_str
                chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
                chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
                chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = COLOR_DARK_BLUE
            else:
                chart.has_title = False

            # Légende
            chart.has_legend = len(chart_spec.get("series", [])) > 1
            if chart.has_legend:
                from pptx.enum.chart import XL_LEGEND_POSITION
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False

            # Couleurs des séries
            for i, series_spec in enumerate(chart_spec.get("series", [])):
                color_hex = series_spec.get("color")
                if color_hex and i < len(chart.series):
                    try:
                        chart.series[i].format.fill.solid()
                        chart.series[i].format.fill.fore_color.rgb = self._hex_to_rgb(color_hex)
                    except Exception:
                        pass

            # Supprime le fond gris par défaut
            chart.plot_area.fill.background()
            chart.chart_area.fill.background()

            return True
        except Exception:
            return False

    def _add_header(self, slide, titre: str, slide_number: int, total_slides: int) -> None:
        """Bande header commune à tous les slides."""
        header = slide.shapes.add_shape(
            1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.3)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_DARK_BLUE
        header.line.fill.background()

        accent_line = slide.shapes.add_shape(
            1, Inches(0), Inches(1.3), SLIDE_WIDTH, Inches(0.05)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = COLOR_ACCENT
        accent_line.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.2), Inches(11.5), Inches(0.9)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = titre
        run.font.size = Pt(22)
        run.font.color.rgb = COLOR_WHITE
        run.font.bold = True
        run.font.name = "Calibri"

        num_box = slide.shapes.add_textbox(
            Inches(12.0), Inches(0.3), Inches(1.1), Inches(0.6)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"{slide_number}/{total_slides}"
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)
        run.font.name = "Calibri"

    def _add_footer(self, slide, sources: list) -> None:
        """Footer sources + ligne accent bas."""
        if sources:
            src_text = "Sources: " + " | ".join(str(s) for s in sources[:4])
            src_box = slide.shapes.add_textbox(
                Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4)
            )
            tf = src_box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = src_text
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(0x95, 0xA5, 0xA6)
            run.font.name = "Calibri"
            run.font.italic = True

        bottom_line = slide.shapes.add_shape(
            1, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.05)
        )
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = COLOR_ACCENT
        bottom_line.line.fill.background()

    def _add_content_slide(
        self,
        prs: Presentation,
        slide_number: int,
        total_slides: int,
        titre: str,
        contenu: str,
        chiffres_cites: list,
        sources: list,
        slide_id: str = "",
        chart_data: dict = None,
    ) -> None:
        """
        Add a content slide.
        - Si chart_data fourni → layout 2 colonnes : texte gauche + graphique droite
        - Sinon → layout original : texte gauche + chiffres clés droite
        """
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        self._add_header(slide, titre, slide_number, total_slides)

        # ── Layout avec graphique ──────────────────────────────────────────
        if chart_data:
            # Texte : colonne gauche (0.5" → 5.5", hauteur 5.0")
            body_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.55), Inches(5.0), Inches(5.0)
            )
            tf = body_box.text_frame
            tf.word_wrap = True
            paragraphs = [p.strip() for p in contenu.split("\n") if p.strip()]
            for i, para_text in enumerate(paragraphs):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_before = Pt(7)
                run = p.add_run()
                # Bullets
                bullet = "• " if not para_text.startswith(("•", "-", "–")) else ""
                run.text = bullet + para_text.lstrip("•-– ")
                run.font.size = Pt(13)
                run.font.color.rgb = COLOR_TEXT
                run.font.name = "Calibri"

            # Chiffres clés compacts sous le texte (si disponibles)
            if chiffres_cites:
                kf_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(6.55), Inches(5.0), Inches(0.55)
                )
                tf = kf_box.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                formatted = []
                for fig in chiffres_cites[:5]:
                    try:
                        n = float(fig)
                        formatted.append(f"{n:,.0f}".replace(",", " ") if n >= 1000 else str(fig))
                    except (ValueError, TypeError):
                        formatted.append(str(fig))
                run.text = "  |  ".join(formatted)
                run.font.size = Pt(10)
                run.font.color.rgb = COLOR_ACCENT
                run.font.bold = True
                run.font.name = "Calibri"

            # Graphique : colonne droite (5.8" → 13.33", hauteur 5.1")
            self._add_chart_to_slide(
                slide, chart_data,
                x=5.8, y=1.55, w=7.2, h=5.1,
            )

        # ── Layout texte + chiffres clés (sans graphique) ─────────────────
        else:
            body_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.55), Inches(8.5), Inches(5.0)
            )
            tf = body_box.text_frame
            tf.word_wrap = True
            paragraphs = [p.strip() for p in contenu.split("\n") if p.strip()]
            for i, para_text in enumerate(paragraphs):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_before = Pt(8)
                run = p.add_run()
                run.text = para_text
                run.font.size = Pt(14)
                run.font.color.rgb = COLOR_TEXT
                run.font.name = "Calibri"

            if chiffres_cites:
                panel = slide.shapes.add_shape(
                    1, Inches(9.3), Inches(1.55), Inches(3.7), Inches(5.0)
                )
                panel.fill.solid()
                panel.fill.fore_color.rgb = COLOR_LIGHT_GREY
                panel.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE0)

                panel_title = slide.shapes.add_textbox(
                    Inches(9.5), Inches(1.75), Inches(3.3), Inches(0.5)
                )
                tf = panel_title.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = "Chiffres clés"
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = COLOR_DARK_BLUE
                run.font.name = "Calibri"

                figures_box = slide.shapes.add_textbox(
                    Inches(9.5), Inches(2.25), Inches(3.3), Inches(4.0)
                )
                tf = figures_box.text_frame
                tf.word_wrap = True
                for i, fig in enumerate(chiffres_cites[:8]):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.space_before = Pt(6)
                    run = p.add_run()
                    try:
                        num = float(fig)
                        if num >= 1_000_000:
                            display = f"{num/1_000_000:.2f} M"
                        elif num >= 1_000:
                            display = f"{num:,.0f}".replace(",", " ")
                        elif num == int(num):
                            display = str(int(num))
                        else:
                            display = f"{num:.2f}"
                        run.text = f"• {display}"
                    except (TypeError, ValueError):
                        run.text = f"• {fig}"
                    run.font.size = Pt(13)
                    run.font.bold = True
                    run.font.color.rgb = COLOR_ACCENT
                    run.font.name = "Calibri"

        self._add_footer(slide, sources)

    def _add_error_slide(
        self, prs: Presentation, slide_id: str, error_message: str
    ) -> None:
        """Add an error placeholder slide when generation failed."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        header = slide.shapes.add_shape(
            1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.3)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_BAD
        header.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"[ERREUR] {slide_id}"
        run.font.size = Pt(20)
        run.font.color.rgb = COLOR_WHITE
        run.font.bold = True
        run.font.name = "Calibri"

        body_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.0)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"Erreur lors de la génération de cette diapositive:\n{error_message}"
        run.font.size = Pt(14)
        run.font.color.rgb = COLOR_TEXT
        run.font.name = "Calibri"

    def build(
        self,
        slides_content: list,
        pharmacy_name: str = "",
    ) -> bytes:
        """
        Build a complete PPTX from generated slide content.

        Args:
            slides_content: List of slide dicts from LLMGenerator.generate_performance_globale().
                            Each dict: {slide_id, titre, contenu, chiffres_cites, sources, erreur}
            pharmacy_name: Name of the pharmacy for the cover slide.

        Returns:
            Raw bytes of the .pptx file, ready for Streamlit download_button.
        """
        prs = self._create_presentation()

        # Add cover slide
        self._add_cover_slide(prs, pharmacy_name)

        # Count valid content slides
        total_content = len(slides_content)

        # Add content slides
        for idx, slide_data in enumerate(slides_content, start=1):
            erreur = slide_data.get("erreur")
            if erreur:
                self._add_error_slide(
                    prs,
                    slide_id=slide_data.get("slide_id", f"slide_{idx}"),
                    error_message=erreur,
                )
            else:
                self._add_content_slide(
                    prs=prs,
                    slide_number=idx,
                    total_slides=total_content,
                    titre=slide_data.get("titre", f"Slide {idx}"),
                    contenu=slide_data.get("contenu", ""),
                    chiffres_cites=slide_data.get("chiffres_cites", []),
                    sources=slide_data.get("sources", []),
                    slide_id=slide_data.get("slide_id", ""),
                    chart_data=slide_data.get("chart_data"),
                )

        # Serialize to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()
