"""
Agence VU — Pipeline de recommandations marketing pharmaceutiques
Streamlit application — Port 8501 (mapped to 19200 externally)
"""

import io
import json
import os
import sys
from pathlib import Path

import pandas as pd
import streamlit as st

# Ensure app directory is on the path
sys.path.insert(0, str(Path(__file__).parent))

from config.settings import get_settings
from parsers.excel_parser import ExcelParser
from parsers.image_parser import ImageParser
from parsers.questionnaire_parser import QuestionnaireParser
from engine.kpi_engine import KPIEngine
from engine.audit import AuditEngine
from engine.data_checker import (
    check_slides_data,
    parse_mapping_from_methodology,
    enrich_kpi_rules,
    build_mapping_markdown,
    DEFAULT_METHODOLOGY_MAPPING,
    STATUS_OK, STATUS_PARTIAL, STATUS_MISSING, STATUS_NO_DATA, STATUS_SKIPPED,
    STATUS_LABELS, STATUS_COLORS,
)
from generation.llm_generator import LLMGenerator
from pptx_builder.assembler import PPTXAssembler
from utils.cost_estimator import estimate_excel_info, estimate_image_cost, format_cost
from utils import project_manager as pm
from utils import methodology_library as ml

# ---------------------------------------------------------------------------
# Page configuration
# ---------------------------------------------------------------------------
st.set_page_config(
    page_title="Agence VU — Pipeline Pharmacie",
    page_icon="💊",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
DATA_DIR = Path("/data")
METHODOLOGY_FILE = DATA_DIR / "methodology.txt"

# ---------------------------------------------------------------------------
# Session state initialization
# ---------------------------------------------------------------------------
def init_session_state():
    defaults = {
        "lot1_excel_results": [],
        "lot1_pptx_texts": {},
        "lot1_image_results": [],
        "lot2_kpis": None,
        "lot2_slides": None,
        "lot2_audit": None,
        "lot2_pptx_bytes": None,
        "lot2_pharmacy_name": "",
        "lot2_context_text": "",
        "lot2_data_check": None,
        "methodology_text": "",
        # Gestion de projet
        "current_project_id": None,
        "current_project_name": None,
        "project_just_loaded": False,
        "confirm_delete": False,
    }
    for key, val in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = val


init_session_state()


# ---------------------------------------------------------------------------
# Professional CSS Theme — Agence VU Design System
# ---------------------------------------------------------------------------
def _inject_css():
    st.markdown(
        """
<style>
/* ── Fonts & Variables ─────────────────────────────────────────────────── */
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');

:root {
  --vu-navy:       #1A2E4A;
  --vu-navy-light: #243d62;
  --vu-blue:       #008BD2;
  --vu-blue-light: #e8f4fd;
  --vu-green:      #27AE60;
  --vu-orange:     #F39C12;
  --vu-red:        #E74C3C;
  --vu-grey:       #95A5A6;
  --text-primary:  #1A2E4A;
  --text-secondary:#5D6D7E;
  --bg-main:       #F0F4F8;
  --bg-card:       #FFFFFF;
  --border:        #DDE3EA;
  --shadow-sm:     0 1px 3px rgba(26,46,74,0.08);
  --shadow-md:     0 4px 12px rgba(26,46,74,0.10);
  --radius:        10px;
  --radius-sm:     6px;
}

/* ── Base ──────────────────────────────────────────────────────────────── */
html, body, [class*="css"] {
  font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif !important;
}

.stApp {
  background: var(--bg-main) !important;
}

.main .block-container {
  padding: 2rem 2.5rem 3rem !important;
  max-width: 1280px;
}

/* ── Sidebar ───────────────────────────────────────────────────────────── */
[data-testid="stSidebar"] {
  background: var(--vu-navy) !important;
  border-right: none !important;
}
[data-testid="stSidebar"] > div {
  background: var(--vu-navy) !important;
}
/* Tous les textes sidebar en blanc */
[data-testid="stSidebar"] label,
[data-testid="stSidebar"] .stMarkdown,
[data-testid="stSidebar"] p,
[data-testid="stSidebar"] span,
[data-testid="stSidebar"] div {
  color: rgba(255,255,255,0.90) !important;
}
[data-testid="stSidebar"] h1,
[data-testid="stSidebar"] h2,
[data-testid="stSidebar"] h3 {
  color: #FFFFFF !important;
}
/* Section labels in sidebar */
[data-testid="stSidebar"] .stMarkdown h3 {
  color: var(--vu-blue) !important;
  font-size: 0.65rem !important;
  font-weight: 600 !important;
  text-transform: uppercase !important;
  letter-spacing: 0.12em !important;
  margin: 1.25rem 0 0.4rem 0 !important;
}
/* Sidebar buttons */
[data-testid="stSidebar"] .stButton > button {
  background: rgba(255,255,255,0.08) !important;
  border: 1px solid rgba(255,255,255,0.18) !important;
  color: rgba(255,255,255,0.9) !important;
  border-radius: var(--radius-sm) !important;
  font-size: 0.8rem !important;
  font-weight: 500 !important;
  transition: all 0.15s ease !important;
}
[data-testid="stSidebar"] .stButton > button:hover {
  background: rgba(0,139,210,0.25) !important;
  border-color: var(--vu-blue) !important;
  color: #FFFFFF !important;
}
/* Sidebar selectbox */
[data-testid="stSidebar"] .stSelectbox [data-baseweb="select"] > div,
[data-testid="stSidebar"] .stSelectbox [data-baseweb="select"] > div:hover {
  background: rgba(255,255,255,0.1) !important;
  border-color: rgba(255,255,255,0.25) !important;
  border-radius: var(--radius-sm) !important;
}
[data-testid="stSidebar"] .stSelectbox svg {
  fill: rgba(255,255,255,0.6) !important;
}
/* Sidebar text input */
[data-testid="stSidebar"] .stTextInput input {
  background: rgba(255,255,255,0.1) !important;
  border-color: rgba(255,255,255,0.25) !important;
  color: white !important;
  border-radius: var(--radius-sm) !important;
}
[data-testid="stSidebar"] .stTextInput input::placeholder {
  color: rgba(255,255,255,0.4) !important;
}
/* Sidebar expander */
[data-testid="stSidebar"] [data-testid="stExpander"] {
  background: rgba(255,255,255,0.06) !important;
  border: 1px solid rgba(255,255,255,0.12) !important;
  border-radius: var(--radius-sm) !important;
}
[data-testid="stSidebar"] [data-testid="stExpander"] summary {
  color: rgba(255,255,255,0.9) !important;
  font-weight: 500 !important;
  font-size: 0.85rem !important;
}
/* Sidebar hr */
[data-testid="stSidebar"] hr {
  border-color: rgba(255,255,255,0.12) !important;
  margin: 0.75rem 0 !important;
}
/* Sidebar alerts */
[data-testid="stSidebar"] [data-testid="stAlert"] {
  background: rgba(255,255,255,0.08) !important;
  border-radius: var(--radius-sm) !important;
  border: 1px solid rgba(255,255,255,0.15) !important;
}
[data-testid="stSidebar"] [data-testid="stAlert"] p {
  font-size: 0.8rem !important;
}
/* Sidebar radio */
[data-testid="stSidebar"] .stRadio > div {
  gap: 0.1rem !important;
}
[data-testid="stSidebar"] .stRadio label {
  font-size: 0.875rem !important;
  font-weight: 500 !important;
  padding: 0.5rem 0.6rem !important;
  border-radius: var(--radius-sm) !important;
  transition: background 0.15s !important;
}
[data-testid="stSidebar"] .stRadio label:hover {
  background: rgba(255,255,255,0.08) !important;
}
/* Active radio option */
[data-testid="stSidebar"] .stRadio [aria-checked="true"] + label,
[data-testid="stSidebar"] .stRadio input:checked + label {
  background: rgba(0,139,210,0.3) !important;
  color: white !important;
}
/* Caption in sidebar */
[data-testid="stSidebar"] .stCaption,
[data-testid="stSidebar"] small {
  color: rgba(255,255,255,0.5) !important;
  font-size: 0.7rem !important;
}

/* ── Main headings ─────────────────────────────────────────────────────── */
.main h1 {
  color: var(--vu-navy) !important;
  font-weight: 700 !important;
  font-size: 1.6rem !important;
  letter-spacing: -0.02em !important;
  padding-bottom: 0.6rem !important;
  border-bottom: 3px solid var(--vu-blue) !important;
  margin-bottom: 1.5rem !important;
}
.main h2 {
  color: var(--vu-navy) !important;
  font-weight: 600 !important;
  font-size: 1.2rem !important;
  margin-top: 1.5rem !important;
  margin-bottom: 0.75rem !important;
}
.main h3 {
  color: var(--vu-navy) !important;
  font-weight: 600 !important;
  font-size: 1rem !important;
}

/* ── Buttons ───────────────────────────────────────────────────────────── */
.stButton > button {
  font-family: 'Inter', sans-serif !important;
  font-weight: 500 !important;
  font-size: 0.875rem !important;
  border-radius: var(--radius-sm) !important;
  transition: all 0.15s ease !important;
  letter-spacing: 0.01em !important;
}
.stButton > button[kind="primary"],
.stButton > button[data-testid*="primary"] {
  background: var(--vu-navy) !important;
  color: white !important;
  border: none !important;
  box-shadow: var(--shadow-sm) !important;
}
.stButton > button[kind="primary"]:hover {
  background: var(--vu-blue) !important;
  box-shadow: 0 4px 14px rgba(0,139,210,0.35) !important;
  transform: translateY(-1px) !important;
}
.stButton > button[kind="secondary"] {
  background: white !important;
  color: var(--vu-navy) !important;
  border: 1px solid var(--border) !important;
}
.stButton > button[kind="secondary"]:hover {
  border-color: var(--vu-blue) !important;
  color: var(--vu-blue) !important;
}

/* ── Metrics ───────────────────────────────────────────────────────────── */
[data-testid="stMetric"] {
  background: var(--bg-card) !important;
  border-radius: var(--radius) !important;
  padding: 1rem 1.25rem !important;
  box-shadow: var(--shadow-sm) !important;
  border: 1px solid var(--border) !important;
  border-left: 4px solid var(--vu-blue) !important;
}
[data-testid="stMetric"] > label {
  color: var(--text-secondary) !important;
  font-size: 0.7rem !important;
  font-weight: 600 !important;
  text-transform: uppercase !important;
  letter-spacing: 0.08em !important;
}
[data-testid="stMetricValue"] {
  color: var(--vu-navy) !important;
  font-weight: 700 !important;
  font-size: 1.5rem !important;
}
[data-testid="stMetricDelta"] {
  font-size: 0.8rem !important;
}

/* ── Expanders / Cards ─────────────────────────────────────────────────── */
[data-testid="stExpander"] {
  background: var(--bg-card) !important;
  border-radius: var(--radius) !important;
  border: 1px solid var(--border) !important;
  box-shadow: var(--shadow-sm) !important;
  margin-bottom: 0.75rem !important;
}
[data-testid="stExpander"] summary {
  font-weight: 600 !important;
  color: var(--vu-navy) !important;
  font-size: 0.9rem !important;
  padding: 0.85rem 1rem !important;
}
[data-testid="stExpander"] summary:hover {
  background: var(--vu-blue-light) !important;
}
[data-testid="stExpander"] summary svg {
  color: var(--vu-blue) !important;
}

/* ── Inputs ────────────────────────────────────────────────────────────── */
.stTextInput input,
.stTextArea textarea,
.stNumberInput input {
  border-radius: var(--radius-sm) !important;
  border-color: var(--border) !important;
  font-family: 'Inter', sans-serif !important;
  font-size: 0.875rem !important;
  color: var(--text-primary) !important;
  background-color: #FFFFFF !important;
  transition: border-color 0.15s, box-shadow 0.15s !important;
}
/* Conteneur wrapper du textarea (Streamlit enveloppe dans un div avec classe css-*) */
.stTextArea > div > div {
  background-color: #FFFFFF !important;
  border-radius: var(--radius-sm) !important;
}
.stTextInput input:focus,
.stTextArea textarea:focus {
  border-color: var(--vu-blue) !important;
  box-shadow: 0 0 0 3px rgba(0,139,210,0.12) !important;
  outline: none !important;
}
.stTextInput label,
.stTextArea label,
.stSelectbox label,
.stNumberInput label {
  color: var(--vu-navy) !important;
  font-weight: 600 !important;
  font-size: 0.8rem !important;
  text-transform: uppercase !important;
  letter-spacing: 0.06em !important;
  margin-bottom: 0.35rem !important;
}

/* ── Selectbox ─────────────────────────────────────────────────────────── */
.stSelectbox [data-baseweb="select"] > div {
  border-radius: var(--radius-sm) !important;
  border-color: var(--border) !important;
  font-family: 'Inter', sans-serif !important;
  font-size: 0.875rem !important;
}
.stSelectbox [data-baseweb="select"] > div:focus-within {
  border-color: var(--vu-blue) !important;
  box-shadow: 0 0 0 3px rgba(0,139,210,0.12) !important;
}

/* ── Alerts ────────────────────────────────────────────────────────────── */
[data-testid="stAlert"] {
  border-radius: var(--radius) !important;
  border: none !important;
  font-size: 0.875rem !important;
}
[data-testid="stAlert"][data-baseweb="notification"][kind="positive"],
.stSuccess {
  background: #EAFAF1 !important;
  border-left: 4px solid var(--vu-green) !important;
  color: #1D6E3B !important;
}
[data-testid="stAlert"][data-baseweb="notification"][kind="warning"],
.stWarning {
  background: #FEF9E7 !important;
  border-left: 4px solid var(--vu-orange) !important;
  color: #9A5C0A !important;
}
[data-testid="stAlert"][data-baseweb="notification"][kind="negative"],
.stError {
  background: #FDEDEC !important;
  border-left: 4px solid var(--vu-red) !important;
  color: #922B21 !important;
}
[data-testid="stAlert"][data-baseweb="notification"][kind="info"],
.stInfo {
  background: var(--vu-blue-light) !important;
  border-left: 4px solid var(--vu-blue) !important;
  color: #1A5276 !important;
}

/* ── File uploader ─────────────────────────────────────────────────────── */
[data-testid="stFileUploader"] {
  background: var(--bg-card) !important;
  border-radius: var(--radius) !important;
  border: 2px dashed var(--border) !important;
  transition: border-color 0.2s !important;
}
[data-testid="stFileUploader"]:hover {
  border-color: var(--vu-blue) !important;
  background: var(--vu-blue-light) !important;
}
[data-testid="stFileUploader"] label {
  font-weight: 600 !important;
  color: var(--vu-navy) !important;
  text-transform: none !important;
  font-size: 0.875rem !important;
  letter-spacing: 0 !important;
}
[data-testid="stFileUploaderDropzone"] {
  border-radius: var(--radius-sm) !important;
}
[data-testid="stFileUploaderDropzoneInstructions"] span {
  color: var(--text-secondary) !important;
  font-size: 0.8rem !important;
}
/* Upload "Browse files" button */
[data-testid="stFileUploaderDropzone"] button {
  background: var(--vu-navy) !important;
  color: white !important;
  border: none !important;
  border-radius: var(--radius-sm) !important;
  font-size: 0.8rem !important;
  font-weight: 600 !important;
  padding: 0.4rem 1rem !important;
  transition: background 0.15s !important;
}
[data-testid="stFileUploaderDropzone"] button:hover {
  background: var(--vu-blue) !important;
}

/* ── DataFrames ────────────────────────────────────────────────────────── */
/* NOTE: pas d'overflow:hidden ni de background sur le contenu interne —
   cela casserait le rendu virtuel de la grille Streamlit */
[data-testid="stDataFrame"] > div:first-child {
  border-radius: var(--radius) !important;
  box-shadow: var(--shadow-sm) !important;
  border: 1px solid var(--border) !important;
}

/* ── Status boxes (st.status) ──────────────────────────────────────────── */
[data-testid="stStatusWidget"] {
  border-radius: var(--radius) !important;
  border: 1px solid var(--border) !important;
  box-shadow: var(--shadow-sm) !important;
  background: var(--bg-card) !important;
}

/* ── Progress bar ──────────────────────────────────────────────────────── */
[data-testid="stProgressBar"] > div > div,
.stProgress > div > div > div {
  background: var(--vu-blue) !important;
  border-radius: 999px !important;
}
[data-testid="stProgressBar"] > div {
  border-radius: 999px !important;
  background: var(--border) !important;
}

/* ── Spinner ───────────────────────────────────────────────────────────── */
[data-testid="stSpinner"] > div {
  border-top-color: var(--vu-blue) !important;
}

/* ── Divider ───────────────────────────────────────────────────────────── */
.main hr {
  border: none !important;
  border-top: 1px solid var(--border) !important;
  margin: 1.5rem 0 !important;
}

/* ── Caption ───────────────────────────────────────────────────────────── */
.stCaption, [data-testid="stCaptionContainer"] p {
  color: var(--text-secondary) !important;
  font-size: 0.78rem !important;
}

/* ── Toast ─────────────────────────────────────────────────────────────── */
[data-testid="stToast"] {
  background: var(--vu-navy) !important;
  color: white !important;
  border-radius: var(--radius) !important;
  box-shadow: var(--shadow-md) !important;
}
[data-testid="stToast"] p {
  color: white !important;
}

/* ── Download button (keep primary style) ─────────────────────────────── */
.stDownloadButton > button {
  border-radius: var(--radius-sm) !important;
  font-weight: 600 !important;
  font-size: 0.875rem !important;
  transition: all 0.15s ease !important;
}
.stDownloadButton > button[kind="primary"] {
  background: var(--vu-navy) !important;
  color: white !important;
  border: none !important;
}
.stDownloadButton > button[kind="primary"]:hover {
  background: var(--vu-blue) !important;
  transform: translateY(-1px) !important;
  box-shadow: 0 4px 14px rgba(0,139,210,0.35) !important;
}

/* ── JSON viewer ───────────────────────────────────────────────────────── */
[data-testid="stJson"] {
  background: var(--bg-card) !important;
  border-radius: var(--radius) !important;
  border: 1px solid var(--border) !important;
}

/* ── Subheader spacing ─────────────────────────────────────────────────── */
.main .stMarkdown h3 {
  margin-top: 1.75rem !important;
  padding-top: 0.5rem !important;
  border-top: 1px solid var(--border) !important;
}
.main .stMarkdown h3:first-child {
  border-top: none !important;
  margin-top: 0 !important;
}

/* ── Column gaps ───────────────────────────────────────────────────────── */
[data-testid="column"] {
  gap: 0.75rem !important;
}

/* ── Scrollbar ─────────────────────────────────────────────────────────── */
::-webkit-scrollbar { width: 6px; height: 6px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: var(--border); border-radius: 999px; }
::-webkit-scrollbar-thumb:hover { background: var(--vu-grey); }

/* ── Texte principal — éviter l'héritage de couleurs inattendues ────────── */
.main .stMarkdown p,
.main .stMarkdown li,
.main .stMarkdown span,
.main p {
  color: var(--text-primary) !important;
}
/* Checklist HTML injectée via unsafe_allow_html */
.main .stMarkdown div[style] {
  color: inherit;
}
/* Force le texte à rester lisible dans les widgets standards */
.stSelectbox [data-baseweb="select"] span,
.stSelectbox [data-baseweb="select"] div[class] {
  color: var(--text-primary) !important;
}
/* st.write / st.text dans les status boxes */
[data-testid="stStatusWidget"] p,
[data-testid="stStatusWidget"] span {
  color: var(--text-primary) !important;
}
</style>
        """,
        unsafe_allow_html=True,
    )


_inject_css()


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def load_methodology() -> str:
    """Load methodology from file if it exists."""
    try:
        if METHODOLOGY_FILE.exists():
            return METHODOLOGY_FILE.read_text(encoding="utf-8")
    except Exception:
        pass
    return ""


def save_methodology(text: str) -> bool:
    """Save methodology text to file."""
    try:
        DATA_DIR.mkdir(parents=True, exist_ok=True)
        METHODOLOGY_FILE.write_text(text, encoding="utf-8")
        return True
    except Exception as exc:
        st.error(f"Erreur lors de la sauvegarde: {exc}")
        return False


def get_api_key() -> str:
    """Get API key from settings or environment."""
    settings = get_settings()
    return settings.anthropic_api_key or os.environ.get("ANTHROPIC_API_KEY", "")


def get_model() -> str:
    """Get model name from settings."""
    settings = get_settings()
    return settings.claude_model or "claude-sonnet-4-6"


def api_key_is_set() -> bool:
    """Check whether an API key is configured."""
    return bool(get_api_key())


def format_kpi_statut(statut: str) -> str:
    """Return a colored emoji for KPI status."""
    return {
        "bon": "🟢 Bon",
        "moyen": "🟡 Moyen",
        "faible": "🔴 Faible",
        "inconnu": "⚪ Inconnu",
    }.get(statut, statut)


def _render_data_checklist(data_check: dict, expanded: bool = True):
    """
    Affiche la checklist de disponibilité des données par slide.
    data_check : résultat de check_slides_data()
    """
    if not data_check:
        return

    slides  = data_check.get("slides", [])
    summary = data_check.get("summary", {})

    # ── Résumé global ──────────────────────────────────────────────────────
    n_ok      = summary.get(STATUS_OK, 0)
    n_partial = summary.get(STATUS_PARTIAL, 0)
    n_missing = summary.get(STATUS_MISSING, 0)
    n_nodata  = summary.get(STATUS_NO_DATA, 0)
    n_skipped = summary.get(STATUS_SKIPPED, 0)

    col_a, col_b, col_c, col_d = st.columns(4)
    col_a.metric("✅ Complets",        n_ok)
    col_b.metric("⚠️ Partiels",       n_partial)
    col_c.metric("❌ Manquants",       n_missing)
    col_d.metric("ℹ️ Sans données",   n_nodata + n_skipped)

    if n_missing == 0 and n_partial == 0:
        st.success("Toutes les données nécessaires ont été trouvées.")
    elif n_missing > 0:
        st.error(
            f"{n_missing} slide(s) sans aucune donnée — les contenus générés seront très génériques."
        )
    else:
        st.warning(
            f"{n_partial} slide(s) avec données partielles — certains chiffres seront manquants."
        )

    # ── Tableau détaillé ───────────────────────────────────────────────────
    with st.expander("📋 Détail par slide", expanded=expanded):
        for slide in slides:
            status    = slide["status"]
            bg_color  = STATUS_COLORS.get(status, "#FFFFFF")
            label     = STATUS_LABELS.get(status, status)

            # En-tête slide — tableau natif plus fiable que HTML injecté
            col_id, col_title, col_status = st.columns([2, 5, 3])
            with col_id:
                st.markdown(
                    f"<p style='margin:0;font-size:0.78rem;font-weight:700;"
                    f"color:#1A2E4A;font-family:monospace;'>{slide['slide_id']}</p>",
                    unsafe_allow_html=True,
                )
            with col_title:
                st.markdown(
                    f"<p style='margin:0;font-size:0.8rem;color:#2C3E50;'>{slide['titre']}</p>",
                    unsafe_allow_html=True,
                )
            with col_status:
                status_html = {
                    STATUS_OK:      "<span style='color:#1D6E3B;font-weight:600;'>✅ Complet</span>",
                    STATUS_PARTIAL: "<span style='color:#9A5C0A;font-weight:600;'>⚠️ Partiel</span>",
                    STATUS_MISSING: "<span style='color:#922B21;font-weight:600;'>❌ Manquant</span>",
                    STATUS_NO_DATA: "<span style='color:#1A5276;'>ℹ️ Sans données</span>",
                    STATUS_SKIPPED: "<span style='color:#7F8C8D;'>⏭️ Non généré</span>",
                }.get(status, label)
                st.markdown(
                    f"<p style='margin:0;font-size:0.8rem;text-align:right;'>{status_html}</p>",
                    unsafe_allow_html=True,
                )

            # Données trouvées
            if slide.get("found_details"):
                for d in slide["found_details"]:
                    if d.get("kpi_id") == "_context":
                        continue
                    st.markdown(
                        f"<p style='margin:0.1rem 0 0 1.5rem;font-size:0.75rem;color:#1D6E3B;'>"
                        f"✔ <strong>{d['label_fr']}</strong> = {d['valeur']} {d['unite']}"
                        f" <em style='color:#7F8C8D;'>(onglet: {d['onglet']})</em></p>",
                        unsafe_allow_html=True,
                    )

            # Données manquantes
            if slide.get("missing_details"):
                for m in slide["missing_details"]:
                    cols_str   = ", ".join(f'&quot;{c}&quot;' for c in m["colonnes_attendues"][:4])
                    sheets_str = ", ".join(f'&quot;{s}&quot;' for s in m["onglets_attendus"][:3])
                    st.markdown(
                        f"<p style='margin:0.1rem 0 0 1.5rem;font-size:0.75rem;color:#922B21;'>"
                        f"✘ <strong>{m['label_fr']}</strong> — non trouvé</p>"
                        f"<p style='margin:0 0 0 2.5rem;font-size:0.72rem;color:#7F8C8D;'>"
                        f"Colonnes : {cols_str or '(non défini)'}  |  "
                        f"Onglets : {sheets_str or '(non défini)'}</p>",
                        unsafe_allow_html=True,
                    )

            # Note spéciale (contexte, no_data, etc.)
            if slide.get("note") and status in (STATUS_NO_DATA, STATUS_SKIPPED):
                st.caption(slide["note"])

            st.markdown(
                "<hr style='margin:0.3rem 0;border:none;border-top:1px solid #EEE;'>",
                unsafe_allow_html=True,
            )


def safe_json_dumps(obj) -> str:
    """Serialize to JSON, handling non-serializable types."""
    def default(o):
        if isinstance(o, float):
            return o
        try:
            return str(o)
        except Exception:
            return None
    return json.dumps(obj, ensure_ascii=False, indent=2, default=default)


# ---------------------------------------------------------------------------
# Sidebar — Projet + Navigation
# ---------------------------------------------------------------------------
with st.sidebar:
    # ── Logo / Branding ──────────────────────────────────────────────────────
    st.markdown(
        """
<div style="
  padding: 1.5rem 1rem 1.25rem;
  margin: -1rem -1rem 0;
  background: linear-gradient(135deg, #1A2E4A 0%, #243d62 100%);
  border-bottom: 2px solid rgba(0,139,210,0.5);
  text-align: center;
">
  <div style="
    font-size: 1.5rem;
    font-weight: 800;
    color: #FFFFFF;
    letter-spacing: 0.1em;
    font-family: 'Inter', sans-serif;
  ">AGENCE <span style="color: #008BD2;">VU</span></div>
  <div style="
    font-size: 0.65rem;
    color: rgba(255,255,255,0.5);
    letter-spacing: 0.15em;
    text-transform: uppercase;
    margin-top: 0.2rem;
  ">Pipeline Pharmacie</div>
</div>
        """,
        unsafe_allow_html=True,
    )

    st.markdown("<div style='height:0.25rem'></div>", unsafe_allow_html=True)

    # ── Projet actif ────────────────────────────────────────────────────────
    st.markdown("### 📁 Projet")

    projects = pm.list_projects()
    project_options = {p["id"]: pm.project_summary(p) for p in projects}

    # Sélecteur de projet existant
    if projects:
        selected_id = st.selectbox(
            "Projet actif",
            options=[""] + list(project_options.keys()),
            format_func=lambda x: "— Sélectionner un projet —" if x == "" else project_options[x],
            index=0 if not st.session_state["current_project_id"]
                  else ([""] + list(project_options.keys())).index(
                      st.session_state["current_project_id"]
                  ) if st.session_state["current_project_id"] in project_options else 0,
            key="sidebar_project_select",
        )
        if selected_id and selected_id != st.session_state["current_project_id"]:
            st.session_state["current_project_id"] = selected_id
            meta = pm.load_project(selected_id)
            st.session_state["current_project_name"] = meta["nom"] if meta else ""
            updated = pm.project_to_session(selected_id, st.session_state)
            st.session_state["project_just_loaded"] = True
            st.rerun()

    # Affichage projet actif
    if st.session_state["current_project_id"]:
        meta = pm.load_project(st.session_state["current_project_id"])
        if meta:
            st.success(f"**{meta['nom']}**\n\n{pm.statut_badge(meta.get('statut','nouveau'))}")

            col_save, col_del = st.columns(2)
            with col_save:
                if st.button("💾 Sauvegarder", use_container_width=True, key="btn_save_project"):
                    pm.session_to_project(st.session_state["current_project_id"], st.session_state)
                    st.toast("✅ Projet sauvegardé !", icon="💾")
            with col_del:
                if st.button("🗑️ Supprimer", use_container_width=True, key="btn_del_project",
                             type="secondary"):
                    st.session_state["confirm_delete"] = True

            if st.session_state.get("confirm_delete"):
                st.warning("Supprimer définitivement ce projet ?")
                c1, c2 = st.columns(2)
                with c1:
                    if st.button("✅ Oui", key="confirm_del_yes"):
                        pm.delete_project(st.session_state["current_project_id"])
                        st.session_state["current_project_id"] = None
                        st.session_state["current_project_name"] = None
                        st.session_state["confirm_delete"] = False
                        st.rerun()
                with c2:
                    if st.button("❌ Non", key="confirm_del_no"):
                        st.session_state["confirm_delete"] = False
                        st.rerun()
    else:
        st.info("Aucun projet actif.")

    # Créer un nouveau projet
    with st.expander("➕ Nouveau projet"):
        new_nom      = st.text_input("Nom du projet", placeholder="Ex: Pharmacie Alésienne 2025", key="new_proj_nom")
        new_pharma   = st.text_input("Nom de la pharmacie", placeholder="Ex: Pharmacie Alésienne", key="new_proj_pharma")
        if st.button("Créer", type="primary", key="btn_create_project"):
            if new_nom and new_pharma:
                meta = pm.create_project(new_nom, new_pharma)
                st.session_state["current_project_id"]   = meta["id"]
                st.session_state["current_project_name"] = meta["nom"]
                st.session_state["lot2_pharmacy_name"]   = meta["pharmacie"]
                # Reset session data pour nouveau projet
                st.session_state["lot1_excel_results"]  = []
                st.session_state["lot1_pptx_texts"]     = {}
                st.session_state["lot1_image_results"]  = []
                st.session_state["lot2_kpis"]           = None
                st.session_state["lot2_slides"]         = None
                st.session_state["lot2_audit"]          = None
                st.session_state["lot2_pptx_bytes"]     = None
                st.session_state["methodology_text"]    = ""
                st.rerun()
            else:
                st.warning("Remplissez les deux champs.")

    st.markdown("---")

    # ── Navigation ──────────────────────────────────────────────────────────
    st.markdown("### Navigation")
    page = st.radio(
        "Sélectionnez une section",
        options=[
            "🔬 Lot 1 — Analyse des exemples",
            "📋 Méthodologie",
            "🚀 Lot 2 — Générer un rapport",
        ],
        label_visibility="collapsed",
    )
    st.markdown("---")

    # ── API status ───────────────────────────────────────────────────────────
    if api_key_is_set():
        st.markdown(
            f"""
<div style="
  margin: 0.25rem 0;
  padding: 0.6rem 0.85rem;
  background: rgba(39,174,96,0.15);
  border: 1px solid rgba(39,174,96,0.35);
  border-radius: 6px;
  display: flex;
  align-items: center;
  gap: 0.5rem;
">
  <span style="color: #4CAF50; font-size: 0.7rem;">●</span>
  <span style="color: rgba(255,255,255,0.85); font-size: 0.78rem; font-weight: 500;">
    API connectée
  </span>
</div>
<div style="
  color: rgba(255,255,255,0.4);
  font-size: 0.68rem;
  padding: 0.2rem 0.1rem;
">Modèle: {get_model()}</div>
            """,
            unsafe_allow_html=True,
        )
    else:
        st.markdown(
            """
<div style="
  margin: 0.25rem 0;
  padding: 0.6rem 0.85rem;
  background: rgba(231,76,60,0.15);
  border: 1px solid rgba(231,76,60,0.35);
  border-radius: 6px;
">
  <span style="color: #E74C3C; font-size: 0.7rem;">●</span>
  <span style="color: rgba(255,255,255,0.85); font-size: 0.78rem; font-weight: 500;">
    Clé API manquante
  </span>
  <div style="color: rgba(255,255,255,0.4); font-size: 0.68rem; margin-top: 0.2rem;">
    Définissez ANTHROPIC_API_KEY dans .env
  </div>
</div>
            """,
            unsafe_allow_html=True,
        )

    st.markdown("<div style='height: 1rem'></div>", unsafe_allow_html=True)
    st.markdown(
        "<div style='color: rgba(255,255,255,0.25); font-size: 0.65rem; "
        "text-align: center; padding-bottom: 0.5rem;'>"
        "Agence VU © 2026 · v1.0</div>",
        unsafe_allow_html=True,
    )


# ===========================================================================
# PAGE 1 — LOT 1: Analyse des exemples
# ===========================================================================
if page == "🔬 Lot 1 — Analyse des exemples":
    st.markdown(
        """
<div style="
  background: linear-gradient(135deg, #1A2E4A 0%, #243d62 100%);
  border-radius: 12px;
  padding: 1.5rem 2rem;
  margin-bottom: 1.75rem;
  display: flex;
  align-items: center;
  gap: 1rem;
">
  <div style="font-size: 2rem;">🔬</div>
  <div>
    <div style="color: white; font-size: 1.35rem; font-weight: 700; letter-spacing: -0.01em;">
      Lot 1 — Analyse des exemples
    </div>
    <div style="color: rgba(255,255,255,0.6); font-size: 0.82rem; margin-top: 0.2rem;">
      Importez vos fichiers exemples pour extraire la structure et la méthodologie des présentations Agence VU
    </div>
  </div>
</div>
        """,
        unsafe_allow_html=True,
    )

    # -----------------------------------------------------------------------
    # Section 0: Charger données existantes (JSON précédemment exporté)
    # -----------------------------------------------------------------------
    with st.expander("📂 Charger des données déjà extraites (JSON)", expanded=True):
        st.markdown(
            "Si vous avez déjà effectué l'extraction et exporté un fichier "
            "`lot1_all_data.json`, chargez-le ici pour restaurer la session "
            "**sans relancer les appels API.**"
        )
        json_restore_file = st.file_uploader(
            "Fichier JSON exporté",
            type=["json"],
            key="lot1_json_restore",
        )
        if json_restore_file:
            try:
                restored = json.loads(json_restore_file.read().decode("utf-8"))
                excel_res  = restored.get("excel", [])
                pptx_res   = restored.get("pptx", {})
                images_res = restored.get("images", [])

                st.session_state["lot1_excel_results"]  = excel_res
                st.session_state["lot1_pptx_texts"]     = pptx_res
                st.session_state["lot1_image_results"]  = images_res

                col_a, col_b, col_c = st.columns(3)
                col_a.metric("Fichiers Excel",   len(excel_res))
                col_b.metric("Slides PPTX",      len(pptx_res))
                col_c.metric("Images extraites", len(images_res))
                st.success("✅ Session restaurée — vous pouvez passer directement à la Méthodologie.")

                # Auto-save si un projet est actif
                pid = st.session_state.get("current_project_id")
                if pid:
                    pm.session_to_project(pid, st.session_state)
                    st.toast("💾 Données sauvegardées dans le projet actif", icon="💾")
            except Exception as exc:
                st.error(f"Impossible de lire le fichier JSON : {exc}")

    # -----------------------------------------------------------------------
    # Section A: Excel files
    # -----------------------------------------------------------------------
    with st.expander("📊 Fichiers Excel (données financières)", expanded=True):
        st.markdown(
            "Importez les fichiers Excel des pharmacies exemples (format `.xlsx`). "
            "Chaque fichier peut contenir plusieurs onglets."
        )
        excel_files = st.file_uploader(
            "Sélectionnez un ou plusieurs fichiers Excel",
            type=["xlsx", "xls"],
            accept_multiple_files=True,
            key="lot1_excel_uploader",
        )

        if excel_files:
            # ── Estimation des coûts avant traitement ──────────────────────
            st.markdown("**📋 Récapitulatif avant traitement**")
            cost_rows = []
            file_bytes_cache = {}
            for uf in excel_files:
                raw = uf.read()
                file_bytes_cache[uf.name] = raw
                info = estimate_excel_info(raw, uf.name)
                cost_rows.append({
                    "Fichier": uf.name,
                    "Taille": f"{info['size_kb']:.1f} Ko",
                    "Onglets": ", ".join(info["sheets"]) if info["sheets"] else "—",
                    "Appel API": "❌ Non (traitement local)",
                    "Coût estimé": "$0.00",
                })
            st.dataframe(
                pd.DataFrame(cost_rows),
                hide_index=True,
                use_container_width=True,
            )
            st.info("✅ Les fichiers Excel sont parsés localement — aucun coût API.")

            parser = ExcelParser()
            all_results = []

            for uploaded_file in excel_files:
                # Réutilise les bytes déjà lus
                uploaded_file._buffer = io.BytesIO(file_bytes_cache[uploaded_file.name])
                with st.spinner(f"Analyse de {uploaded_file.name}..."):
                    try:
                        file_bytes = io.BytesIO(file_bytes_cache[uploaded_file.name])
                        file_bytes.name = uploaded_file.name
                        result = parser.parse(file_bytes)
                        all_results.append(result)

                        st.subheader(f"📄 {uploaded_file.name}")

                        # Summary table
                        sheet_summary = []
                        for sname, sdata in result["sheets"].items():
                            sheet_summary.append(
                                {
                                    "Onglet": sname,
                                    "Colonnes": len(sdata["headers"]),
                                    "Lignes de données": len(sdata["rows"]),
                                    "Cellules numériques": len(sdata["numeric_cells"]),
                                }
                            )
                        if sheet_summary:
                            st.dataframe(
                                pd.DataFrame(sheet_summary),
                                use_container_width=True,
                                hide_index=True,
                            )

                        # Detailed JSON (collapsed)
                        with st.expander(f"Données brutes — {uploaded_file.name}"):
                            st.json(result)

                    except Exception as exc:
                        st.error(f"Erreur lors du parsing de {uploaded_file.name}: {exc}")

            if all_results:
                st.session_state["lot1_excel_results"] = all_results

                # Export button
                export_data = safe_json_dumps(all_results)
                st.download_button(
                    label="💾 Exporter données Excel extraites (JSON)",
                    data=export_data.encode("utf-8"),
                    file_name="lot1_excel_data.json",
                    mime="application/json",
                )

    # -----------------------------------------------------------------------
    # Section B: PPTX example
    # -----------------------------------------------------------------------
    with st.expander("📑 Fichier PPTX exemple", expanded=False):
        st.markdown(
            "Importez un exemple de présentation PPTX Agence VU pour analyser "
            "la structure et le contenu attendus."
        )
        pptx_file = st.file_uploader(
            "Sélectionnez un fichier PowerPoint",
            type=["pptx"],
            accept_multiple_files=False,
            key="lot1_pptx_uploader",
        )

        if pptx_file:
            with st.spinner(f"Analyse de {pptx_file.name}..."):
                try:
                    from pptx import Presentation as PPTXPresentation

                    prs = PPTXPresentation(io.BytesIO(pptx_file.read()))
                    slide_texts = {}

                    for i, slide in enumerate(prs.slides, start=1):
                        texts = []
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    line = para.text.strip()
                                    if line:
                                        texts.append(line)
                        if texts:
                            slide_key = f"Diapositive {i}"
                            slide_texts[slide_key] = texts

                    st.session_state["lot1_pptx_texts"] = slide_texts

                    st.success(
                        f"✅ {len(prs.slides)} diapositives extraites de {pptx_file.name}"
                    )

                    # Display per slide
                    for slide_label, lines in slide_texts.items():
                        with st.expander(slide_label):
                            for line in lines:
                                st.markdown(f"- {line}")

                    # Export
                    export_data = safe_json_dumps(slide_texts)
                    st.download_button(
                        label="💾 Exporter textes PPTX (JSON)",
                        data=export_data.encode("utf-8"),
                        file_name="lot1_pptx_texts.json",
                        mime="application/json",
                    )

                except Exception as exc:
                    st.error(f"Erreur lors du parsing PPTX: {exc}")

    # -----------------------------------------------------------------------
    # Section C: Images PNG/JPG
    # -----------------------------------------------------------------------
    with st.expander("🖼️ Images PNG (exports tableaux de bord)", expanded=False):
        st.markdown(
            "Importez des captures d'écran ou exports PNG/JPG des logiciels de gestion. "
            "Claude Vision extraira automatiquement les valeurs numériques visibles."
        )

        if not api_key_is_set():
            st.warning(
                "⚠️ Clé API Anthropic non configurée. "
                "Définissez `ANTHROPIC_API_KEY` dans le fichier `.env` pour utiliser "
                "la vision Claude."
            )

        image_files = st.file_uploader(
            "Sélectionnez une ou plusieurs images",
            type=["png", "jpg", "jpeg"],
            accept_multiple_files=True,
            key="lot1_image_uploader",
        )

        if image_files:
            # ── Estimation des coûts avant traitement ──────────────────────
            st.markdown("**📋 Récapitulatif avant traitement**")
            img_cost_rows = []
            img_bytes_cache = {}
            total_cost = 0.0
            for img_file in image_files:
                raw = img_file.read()
                img_bytes_cache[img_file.name] = raw
                est = estimate_image_cost(raw)
                total_cost += est["total_cost_usd"]
                img_cost_rows.append({
                    "Fichier": img_file.name,
                    "Taille": f"{len(raw)/1024:.1f} Ko",
                    "Dimensions": f"{est['width']}×{est['height']} px",
                    "Tokens entrée": f"{est['input_tokens']:,}",
                    "Tokens sortie": f"~{est['output_tokens']:,}",
                    "Coût estimé": format_cost(est["total_cost_usd"]),
                })
            st.dataframe(
                pd.DataFrame(img_cost_rows),
                hide_index=True,
                use_container_width=True,
            )
            st.info(
                f"💰 **Coût total estimé pour {len(image_files)} image(s) : "
                f"{format_cost(total_cost)}**  \n"
                f"_(entrée $3/MTok + sortie $15/MTok — modèle {get_model()})_"
            )

            if not api_key_is_set():
                st.error(
                    "Impossible d'analyser les images sans clé API Anthropic. "
                    "Ajoutez votre clé dans le fichier `.env`."
                )
            else:
                image_parser = ImageParser(
                    api_key=get_api_key(),
                    model=get_model(),
                )
                all_image_results = []

                for img_file in image_files:
                    with st.spinner(f"Analyse de {img_file.name} avec Claude Vision..."):
                        try:
                            img_bytes = img_bytes_cache[img_file.name]
                            result = image_parser.parse(img_bytes, filename=img_file.name)
                            all_image_results.append(result)

                            st.subheader(f"🖼️ {img_file.name}")

                            if result.get("erreur"):
                                st.error(f"Erreur: {result['erreur']}")
                            else:
                                valeurs = result.get("valeurs_extraites", [])
                                if valeurs:
                                    df_vals = pd.DataFrame(valeurs)
                                    # Add color styling for confiance
                                    st.dataframe(
                                        df_vals,
                                        use_container_width=True,
                                        hide_index=True,
                                    )
                                    st.caption(
                                        f"{len(valeurs)} valeur(s) extraite(s)"
                                    )
                                else:
                                    st.info("Aucune valeur numérique extraite.")

                        except Exception as exc:
                            st.error(f"Erreur lors du traitement de {img_file.name}: {exc}")

                if all_image_results:
                    st.session_state["lot1_image_results"] = all_image_results

                    export_data = safe_json_dumps(all_image_results)
                    st.download_button(
                        label="💾 Exporter données images extraites (JSON)",
                        data=export_data.encode("utf-8"),
                        file_name="lot1_image_data.json",
                        mime="application/json",
                    )

    # -----------------------------------------------------------------------
    # Global export
    # -----------------------------------------------------------------------
    st.markdown("---")
    if (
        st.session_state["lot1_excel_results"]
        or st.session_state["lot1_pptx_texts"]
        or st.session_state["lot1_image_results"]
    ):
        st.markdown(
            "<div style='font-size:0.65rem;font-weight:700;color:#5D6D7E;text-transform:uppercase;"
            "letter-spacing:0.12em;margin:1.5rem 0 0.75rem;'>💾 &nbsp;Export global</div>",
            unsafe_allow_html=True,
        )
        all_lot1_data = {
            "excel": st.session_state["lot1_excel_results"],
            "pptx": st.session_state["lot1_pptx_texts"],
            "images": st.session_state["lot1_image_results"],
        }
        export_all = safe_json_dumps(all_lot1_data)
        st.download_button(
            label="💾 Exporter TOUTES les données extraites (JSON)",
            data=export_all.encode("utf-8"),
            file_name="lot1_all_data.json",
            mime="application/json",
            type="primary",
        )


# ===========================================================================
# PAGE 2 — Méthodologie
# ===========================================================================
elif page == "📋 Méthodologie":
    st.markdown(
        """
<div style="
  background: linear-gradient(135deg, #1A2E4A 0%, #243d62 100%);
  border-radius: 12px;
  padding: 1.5rem 2rem;
  margin-bottom: 1.75rem;
  display: flex;
  align-items: center;
  gap: 1rem;
">
  <div style="font-size: 2rem;">📋</div>
  <div>
    <div style="color: white; font-size: 1.35rem; font-weight: 700; letter-spacing: -0.01em;">
      Méthodologie Agence VU
    </div>
    <div style="color: rgba(255,255,255,0.6); font-size: 0.82rem; margin-top: 0.2rem;">
      Gérez, générez et appliquez les méthodologies de génération de rapports
    </div>
  </div>
</div>
        """,
        unsafe_allow_html=True,
    )

    # Clés session pour la page
    if "methodo_active_id"   not in st.session_state: st.session_state["methodo_active_id"]   = None
    if "methodo_rename_id"   not in st.session_state: st.session_state["methodo_rename_id"]   = None
    if "methodo_confirm_del" not in st.session_state: st.session_state["methodo_confirm_del"] = None

    # Chargement initial depuis fichier global (rétrocompat)
    if not st.session_state["methodology_text"]:
        st.session_state["methodology_text"] = load_methodology()

    methodo_list = ml.list_methodologies()

    # ===========================================================================
    # BLOC 1 — Bibliothèque de méthodologies
    # ===========================================================================
    st.markdown("### 📚 Bibliothèque de méthodologies")

    if not methodo_list:
        st.info("Aucune méthodologie sauvegardée. Générez-en une depuis le PPTX exemple ci-dessous.")
    else:
        # ── Sélecteur ────────────────────────────────────────────────────────
        options_ids   = [m["id"]  for m in methodo_list]
        options_labels= [ml.summary_line(m) for m in methodo_list]

        current_idx = 0
        if st.session_state["methodo_active_id"] in options_ids:
            current_idx = options_ids.index(st.session_state["methodo_active_id"])

        sel_idx = st.selectbox(
            "Sélectionner une méthodologie",
            options=range(len(options_ids)),
            format_func=lambda i: options_labels[i],
            index=current_idx,
            key="methodo_selector",
        )
        sel_id = options_ids[sel_idx]

        col_load, col_dup, col_dl, _ = st.columns([2, 2, 2, 3])
        with col_load:
            if st.button("📂 Charger dans l'éditeur", type="primary", use_container_width=True):
                content = ml.get_content(sel_id)
                st.session_state["methodology_text"]  = content
                st.session_state["methodo_active_id"] = sel_id
                st.rerun()
        with col_dup:
            if st.button("📋 Dupliquer", use_container_width=True):
                src = next((m for m in methodo_list if m["id"] == sel_id), None)
                if src:
                    ml.duplicate_methodology(sel_id, f"Copie — {src['nom']}")
                    st.rerun()
        with col_dl:
            content_dl = ml.get_content(sel_id)
            nom_dl = next((m["nom"] for m in methodo_list if m["id"] == sel_id), sel_id)
            st.download_button(
                "📥 Télécharger",
                data=content_dl.encode("utf-8"),
                file_name=f"{sel_id}.md",
                mime="text/markdown",
                use_container_width=True,
            )

        st.markdown("---")

        # ── Tableau de gestion ───────────────────────────────────────────────
        with st.expander("🗂️ Gérer toutes les méthodologies", expanded=False):
            for m in methodo_list:
                is_active = (m["id"] == st.session_state.get("methodo_active_id"))
                badge = " 🟢 **[active]**" if is_active else ""
                st.markdown(f"**{m['nom']}**{badge} — {m.get('word_count',0)} mots — modifié le {m.get('date_modification','')[:10]}")
                if m.get("description"):
                    st.caption(m["description"])
                if m.get("source_pptx"):
                    st.caption(f"Source PPTX : {m['source_pptx']}")

                c1, c2, c3, c4 = st.columns([2, 2, 2, 5])
                with c1:
                    if st.button("📂 Charger", key=f"load_{m['id']}", use_container_width=True):
                        content = ml.get_content(m["id"])
                        st.session_state["methodology_text"]  = content
                        st.session_state["methodo_active_id"] = m["id"]
                        st.rerun()
                with c2:
                    if st.button("✏️ Renommer", key=f"ren_{m['id']}", use_container_width=True):
                        st.session_state["methodo_rename_id"] = m["id"]
                with c3:
                    if st.button("🗑️ Supprimer", key=f"del_{m['id']}", use_container_width=True):
                        st.session_state["methodo_confirm_del"] = m["id"]

                # Renommer inline
                if st.session_state["methodo_rename_id"] == m["id"]:
                    new_name = st.text_input("Nouveau nom", value=m["nom"], key=f"rename_input_{m['id']}")
                    rc1, rc2 = st.columns(2)
                    with rc1:
                        if st.button("✅ Valider", key=f"rename_ok_{m['id']}"):
                            ml.rename_methodology(m["id"], new_name)
                            st.session_state["methodo_rename_id"] = None
                            st.rerun()
                    with rc2:
                        if st.button("❌ Annuler", key=f"rename_cancel_{m['id']}"):
                            st.session_state["methodo_rename_id"] = None
                            st.rerun()

                # Confirmation suppression
                if st.session_state["methodo_confirm_del"] == m["id"]:
                    st.warning(f"Supprimer définitivement « {m['nom']} » ?")
                    dc1, dc2 = st.columns(2)
                    with dc1:
                        if st.button("✅ Oui, supprimer", key=f"del_ok_{m['id']}"):
                            ml.delete_methodology(m["id"])
                            if st.session_state["methodo_active_id"] == m["id"]:
                                st.session_state["methodo_active_id"] = None
                                st.session_state["methodology_text"]  = ""
                            st.session_state["methodo_confirm_del"] = None
                            st.rerun()
                    with dc2:
                        if st.button("❌ Annuler", key=f"del_cancel_{m['id']}"):
                            st.session_state["methodo_confirm_del"] = None
                            st.rerun()

                st.markdown("---")

    # ── Import externe ────────────────────────────────────────────────────────
    with st.expander("📤 Importer une méthodologie (.md / .txt)"):
        imp_col1, imp_col2 = st.columns([3, 2])
        with imp_col1:
            methodo_upload = st.file_uploader(
                "Fichier Markdown ou texte",
                type=["md", "txt"],
                key="methodo_import_uploader",
                label_visibility="collapsed",
            )
        with imp_col2:
            import_nom = st.text_input("Nom à donner", placeholder="Ex: Benchmark 2025", key="import_nom")

        if methodo_upload and st.button("📥 Importer et sauvegarder", type="primary"):
            imported = methodo_upload.read().decode("utf-8")
            nom = import_nom.strip() or methodo_upload.name.rsplit(".", 1)[0]
            meta = ml.save_methodology(nom=nom, content=imported, source_pptx="import manuel")
            st.session_state["methodology_text"]  = imported
            st.session_state["methodo_active_id"] = meta["id"]
            save_methodology(imported)   # sync fichier global
            st.success(f"✅ Importée et sauvegardée sous « {nom} »")
            st.rerun()

    st.markdown("---")

    # ===========================================================================
    # BLOC 2 — Génération depuis PPTX
    # ===========================================================================
    st.markdown("### 🤖 Générer depuis un PPTX exemple")

    pptx_texts = st.session_state.get("lot1_pptx_texts", {})
    n_slides   = len(pptx_texts)

    gen_col1, gen_col2 = st.columns([2, 1])

    with gen_col1:
        if n_slides:
            st.success(f"✅ PPTX exemple chargé — **{n_slides} diapositives** disponibles")
        else:
            st.warning("⚠️ Aucun PPTX chargé. Allez dans **Lot 1 → Section PPTX** pour en importer un.")

        gen_nom = st.text_input(
            "Nom de la nouvelle méthodologie",
            placeholder="Ex: Pharmacie urbaine sous-pivot 2025",
            key="gen_methodo_nom",
        )

    with gen_col2:
        if not api_key_is_set():
            st.error("Clé API manquante")
        elif n_slides > 0:
            from generation.methodology_generator import MethodologyGenerator
            _mg_tmp = MethodologyGenerator(api_key=get_api_key(), model=get_model())
            est = _mg_tmp.estimate_cost(pptx_texts)
            st.caption(
                f"~{est['slides_selected']}/{est['slides_total']} slides sélectionnées  \n"
                f"~{est['input_tokens']:,} tokens  \n"
                f"Coût estimé : **${est['total_cost_usd']:.4f}**"
            )

    if n_slides > 0 and api_key_is_set():
        if st.button("🤖 Générer la méthodologie", type="primary", use_container_width=True):
            nom_to_use = gen_nom.strip() or "Méthodologie générée"
            with st.spinner("Claude analyse le PPTX et induit la méthodologie…"):
                try:
                    from generation.methodology_generator import MethodologyGenerator
                    mg = MethodologyGenerator(api_key=get_api_key(), model=get_model())
                    generated = mg.generate(pptx_texts)

                    # Sauvegarde dans la bibliothèque
                    pptx_source = "PPTX chargé en Lot 1"
                    meta = ml.save_methodology(
                        nom=nom_to_use,
                        content=generated,
                        source_pptx=pptx_source,
                    )
                    # Charge dans l'éditeur + projet
                    st.session_state["methodology_text"]  = generated
                    st.session_state["methodo_active_id"] = meta["id"]
                    save_methodology(generated)
                    pid = st.session_state.get("current_project_id")
                    if pid:
                        pm.save_methodology(pid, generated)

                    st.success(f"✅ Méthodologie « {nom_to_use} » générée et sauvegardée dans la bibliothèque !")
                    st.rerun()
                except Exception as exc:
                    st.error(f"Erreur : {exc}")

    st.markdown("---")

    # ===========================================================================
    # BLOC 3 — Éditeur
    # ===========================================================================
    active_nom = ""
    if st.session_state.get("methodo_active_id"):
        active_meta = next((m for m in methodo_list if m["id"] == st.session_state["methodo_active_id"]), None)
        if active_meta:
            active_nom = active_meta["nom"]

    st.markdown(f"### ✏️ Éditeur{f' — *{active_nom}*' if active_nom else ''}")

    methodology_content = st.text_area(
        "Contenu de la méthodologie",
        value=st.session_state["methodology_text"],
        height=550,
        placeholder="Chargez une méthodologie depuis la bibliothèque ou générez-en une depuis le PPTX exemple.",
        key="methodology_textarea",
    )

    if methodology_content:
        st.caption(f"📝 {len(methodology_content.split())} mots — {len(methodology_content)} caractères")

    # Sauvegarde
    save_cols = st.columns([3, 3, 3])

    with save_cols[0]:
        # Écrase la méthodologie active
        overwrite_disabled = not (methodology_content and st.session_state.get("methodo_active_id"))
        if st.button(
            f"💾 Mettre à jour « {active_nom} »" if active_nom else "💾 Mettre à jour",
            type="primary",
            disabled=overwrite_disabled,
            use_container_width=True,
        ):
            ml.save_methodology(
                nom=active_nom,
                content=methodology_content,
                methodo_id=st.session_state["methodo_active_id"],
            )
            st.session_state["methodology_text"] = methodology_content
            save_methodology(methodology_content)
            pid = st.session_state.get("current_project_id")
            if pid:
                pm.save_methodology(pid, methodology_content)
            st.success("✅ Mise à jour sauvegardée")

    with save_cols[1]:
        # Sauvegarde sous un nouveau nom
        new_save_nom = st.text_input("", placeholder="Nouveau nom…", key="save_as_nom", label_visibility="collapsed")
        if st.button("💾 Sauvegarder sous…", use_container_width=True, disabled=not methodology_content):
            nom_final = new_save_nom.strip() or "Sans titre"
            meta = ml.save_methodology(nom=nom_final, content=methodology_content)
            st.session_state["methodo_active_id"] = meta["id"]
            save_methodology(methodology_content)
            st.success(f"✅ Sauvegardée sous « {nom_final} »")
            st.rerun()

    with save_cols[2]:
        if methodology_content:
            fname = f"{st.session_state.get('methodo_active_id', 'methodologie')}.md"
            st.download_button(
                "📥 Télécharger (.md)",
                data=methodology_content.encode("utf-8"),
                file_name=fname,
                mime="text/markdown",
                use_container_width=True,
            )

    # ── Section mapping données ───────────────────────────────────────────────
    if methodology_content:
        has_mapping = "<!-- VU_MAPPING -->" in methodology_content
        with st.expander(
            "📐 Mapping des données" + (" ✅ Présent" if has_mapping else " ⚠️ Absent"),
            expanded=not has_mapping,
        ):
            if has_mapping:
                parsed = parse_mapping_from_methodology(methodology_content)
                st.success(
                    f"Section de mapping présente — **{len(parsed)} KPI(s)** mappés. "
                    "Le pipeline utilisera ces noms de colonnes en priorité."
                )
                st.json(parsed)
            else:
                st.warning(
                    "Cette méthodologie ne contient pas de section de mapping des données. "
                    "Sans elle, le pipeline utilise des noms de colonnes génériques et risque "
                    "de ne pas retrouver vos données Excel."
                )
                if st.button(
                    "➕ Insérer le mapping par défaut dans la méthodologie",
                    type="primary",
                ):
                    mapping_block = build_mapping_markdown(DEFAULT_METHODOLOGY_MAPPING)
                    new_content = methodology_content + mapping_block
                    st.session_state["methodology_text"] = new_content
                    # Sauvegarde automatique si méthodologie active
                    if st.session_state.get("methodo_active_id"):
                        active_meta = next(
                            (m for m in ml.list_methodologies()
                             if m["id"] == st.session_state["methodo_active_id"]),
                            None,
                        )
                        if active_meta:
                            ml.save_methodology(
                                nom=active_meta["nom"],
                                content=new_content,
                                methodo_id=st.session_state["methodo_active_id"],
                            )
                    st.success("✅ Section mapping insérée et sauvegardée. Éditez les noms de colonnes selon vos fichiers Excel.")
                    st.rerun()

    # Aperçu
    if methodology_content:
        with st.expander("👁️ Aperçu rendu"):
            st.markdown(methodology_content)


# ===========================================================================
# PAGE 3 — LOT 2: Générer un rapport
# ===========================================================================
elif page == "🚀 Lot 2 — Générer un rapport":
    st.markdown(
        """
<div style="
  background: linear-gradient(135deg, #1A2E4A 0%, #243d62 100%);
  border-radius: 12px;
  padding: 1.5rem 2rem;
  margin-bottom: 1.75rem;
  display: flex;
  align-items: center;
  gap: 1rem;
">
  <div style="font-size: 2rem;">🚀</div>
  <div>
    <div style="color: white; font-size: 1.35rem; font-weight: 700; letter-spacing: -0.01em;">
      Lot 2 — Génération automatique de rapport
    </div>
    <div style="color: rgba(255,255,255,0.6); font-size: 0.82rem; margin-top: 0.2rem;">
      Pipeline complet · Section <strong style="color: #008BD2;">Performance Globale</strong> · AUDIT 360°
    </div>
  </div>
</div>
        """,
        unsafe_allow_html=True,
    )

    # -----------------------------------------------------------------------
    # Configuration
    # -----------------------------------------------------------------------
    st.markdown(
        "<div style='font-size:0.65rem;font-weight:700;color:#5D6D7E;text-transform:uppercase;"
        "letter-spacing:0.12em;margin-bottom:0.75rem;'>⚙️ &nbsp;Configuration</div>",
        unsafe_allow_html=True,
    )
    col1, col2 = st.columns(2)

    with col1:
        pharmacy_name = st.text_input(
            "Nom de la pharmacie",
            value=st.session_state["lot2_pharmacy_name"],
            placeholder="Ex: Pharmacie du Marché",
            help="Ce nom apparaîtra sur la page de couverture du PowerPoint.",
        )
        st.session_state["lot2_pharmacy_name"] = pharmacy_name

    with col2:
        if not api_key_is_set():
            st.error("⚠️ Clé API Anthropic non configurée (nécessaire pour les étapes 3+)")
        else:
            st.success(f"✅ API prête — Modèle: `{get_model()}`")

    st.markdown("---")

    # -----------------------------------------------------------------------
    # File uploads
    # -----------------------------------------------------------------------
    st.markdown(
        "<div style='font-size:0.65rem;font-weight:700;color:#5D6D7E;text-transform:uppercase;"
        "letter-spacing:0.12em;margin:1.5rem 0 0.75rem;'>📁 &nbsp;Fichiers d'entrée</div>",
        unsafe_allow_html=True,
    )

    upload_col1, upload_col2, upload_col3, upload_col4 = st.columns(4)

    with upload_col1:
        st.markdown("**📊 Fichiers Excel**")
        lot2_excel_files = st.file_uploader(
            "Fichiers Excel (.xlsx)",
            type=["xlsx", "xls"],
            accept_multiple_files=True,
            key="lot2_excel_uploader",
        )
        if lot2_excel_files:
            st.caption(f"✅ {len(lot2_excel_files)} fichier(s) chargé(s)")

    with upload_col2:
        st.markdown("**🖼️ Images tableaux de bord**")
        lot2_image_files = st.file_uploader(
            "Images PNG/JPG",
            type=["png", "jpg", "jpeg"],
            accept_multiple_files=True,
            key="lot2_image_uploader",
        )
        if lot2_image_files:
            st.caption(f"✅ {len(lot2_image_files)} image(s) chargée(s)")

    with upload_col3:
        st.markdown("**📋 Questionnaire**")
        lot2_questionnaire_file = st.file_uploader(
            "PDF ou JSON",
            type=["pdf", "json"],
            accept_multiple_files=False,
            key="lot2_questionnaire_uploader",
        )
        if lot2_questionnaire_file:
            ext = lot2_questionnaire_file.name.rsplit(".", 1)[-1].lower()
            if ext == "pdf":
                st.caption("✅ PDF chargé — structuration via Claude")
            else:
                st.caption("✅ JSON chargé")

        with st.expander("Ou saisir manuellement (JSON)"):
            questionnaire_manual = st.text_area(
                "JSON du questionnaire",
                value="{}",
                height=150,
                help="Format JSON: {\"clé\": valeur}",
                key="questionnaire_manual_input",
            )

    with upload_col4:
        st.markdown("**📄 Contexte pharmacie (PDF)**")
        st.caption("Génère le 1er slide d'introduction")
        lot2_context_pdf = st.file_uploader(
            "Document de contexte PDF",
            type=["pdf"],
            accept_multiple_files=False,
            key="lot2_context_uploader",
        )
        if lot2_context_pdf:
            st.caption(f"✅ {lot2_context_pdf.name}")
        elif st.session_state.get("lot2_context_text"):
            st.caption("✅ Contexte déjà extrait en session")

    st.markdown("---")

    # ── Estimation des coûts Lot 2 ──────────────────────────────────────────
    if lot2_excel_files or lot2_image_files:
        st.markdown(
        "<div style='font-size:0.65rem;font-weight:700;color:#5D6D7E;text-transform:uppercase;"
        "letter-spacing:0.12em;margin:1.5rem 0 0.75rem;'>💰 &nbsp;Estimation des coûts API</div>",
        unsafe_allow_html=True,
    )
        preview_rows = []

        for uf in (lot2_excel_files or []):
            preview_rows.append({
                "Fichier": uf.name,
                "Type": "Excel",
                "Taille": f"{uf.size/1024:.1f} Ko",
                "Appel API": "❌ Non",
                "Coût estimé": "$0.00",
            })

        lot2_img_total = 0.0
        for img_file in (lot2_image_files or []):
            raw = img_file.getvalue()  # BytesIO-like, ne consomme pas le curseur
            est = estimate_image_cost(raw)
            lot2_img_total += est["total_cost_usd"]
            preview_rows.append({
                "Fichier": img_file.name,
                "Type": "Image (Vision)",
                "Taille": f"{len(raw)/1024:.1f} Ko",
                "Appel API": "✅ Oui",
                "Coût estimé": format_cost(est["total_cost_usd"]),
            })

        if preview_rows:
            st.dataframe(
                pd.DataFrame(preview_rows),
                hide_index=True,
                use_container_width=True,
            )
            if lot2_img_total > 0:
                st.info(
                    f"💰 **Coût images estimé : {format_cost(lot2_img_total)}** "
                    f"+ coût génération LLM (~{format_cost(0.015)} par rapport)  \n"
                    f"_(modèle {get_model()})_"
                )

    # -----------------------------------------------------------------------
    # Launch pipeline
    # -----------------------------------------------------------------------
    if not pharmacy_name:
        st.markdown(
            """
<div style="
  background: #FEF9E7;
  border-left: 4px solid #F39C12;
  border-radius: 8px;
  padding: 0.75rem 1rem;
  margin: 0.5rem 0;
  font-size: 0.85rem;
  color: #9A5C0A;
">
  ⚠️ &nbsp;Veuillez saisir le nom de la pharmacie avant de lancer le pipeline.
</div>
            """,
            unsafe_allow_html=True,
        )

    launch_disabled = not pharmacy_name

    st.markdown("<div style='height:0.5rem'></div>", unsafe_allow_html=True)
    if st.button(
        "▶  Lancer le pipeline complet",
        type="primary",
        disabled=launch_disabled,
        use_container_width=True,
    ):
        # Reset previous results
        st.session_state["lot2_kpis"]         = None
        st.session_state["lot2_slides"]        = None
        st.session_state["lot2_audit"]         = None
        st.session_state["lot2_pptx_bytes"]    = None
        st.session_state["lot2_context_text"]  = ""
        st.session_state["lot2_data_check"]    = None

        # ===================================================================
        # STEP 1 — Parse all inputs
        # ===================================================================
        with st.status("Étape 1 — Parsing des fichiers d'entrée...", expanded=True) as status:
            st.write("Initialisation des parseurs...")

            excel_parser = ExcelParser()
            all_excel_data = {"sheets": {}, "source": "combined", "total_sheets": 0}
            parse_errors = []

            # Parse Excel files
            if lot2_excel_files:
                for uploaded_file in lot2_excel_files:
                    st.write(f"  📊 Parsing Excel: {uploaded_file.name}")
                    try:
                        file_bytes = io.BytesIO(uploaded_file.read())
                        file_bytes.name = uploaded_file.name
                        result = excel_parser.parse(file_bytes)
                        # Merge sheets
                        for sname, sdata in result["sheets"].items():
                            key = f"{uploaded_file.name}::{sname}"
                            all_excel_data["sheets"][key] = sdata
                        all_excel_data["total_sheets"] += result["total_sheets"]
                        all_excel_data["source"] = uploaded_file.name
                        st.write(f"    ✅ {result['total_sheets']} onglet(s) parsés")
                    except Exception as exc:
                        parse_errors.append(f"Excel {uploaded_file.name}: {exc}")
                        st.write(f"    ❌ Erreur: {exc}")
            else:
                st.write("  ⚠️ Aucun fichier Excel fourni")

            # Parse images
            image_kv_data = {}
            if lot2_image_files and api_key_is_set():
                image_parser = ImageParser(
                    api_key=get_api_key(),
                    model=get_model(),
                )
                for img_file in lot2_image_files:
                    st.write(f"  🖼️ Analyse image: {img_file.name}")
                    try:
                        img_bytes = img_file.read()
                        result = image_parser.parse(img_bytes, filename=img_file.name)
                        if result.get("erreur"):
                            st.write(f"    ❌ {result['erreur']}")
                        else:
                            n = len(result.get("valeurs_extraites", []))
                            st.write(f"    ✅ {n} valeur(s) extraite(s)")
                            # Convert image values to a pseudo-sheet for KPI engine
                            img_sheet_name = f"IMG::{img_file.name}"
                            headers = ["label", "valeur", "unite", "confiance"]
                            rows = []
                            numeric_cells = []
                            for idx, val_entry in enumerate(result.get("valeurs_extraites", []), start=2):
                                v = val_entry.get("valeur")
                                row = [
                                    val_entry.get("label", ""),
                                    v,
                                    val_entry.get("unite", ""),
                                    val_entry.get("confiance", ""),
                                ]
                                rows.append(row)
                                if isinstance(v, (int, float)) and v is not None:
                                    numeric_cells.append({
                                        "ref": f"B{idx}",
                                        "valeur": float(v),
                                        "sheet": img_sheet_name,
                                        "row": idx,
                                        "col": 2,
                                    })
                            all_excel_data["sheets"][img_sheet_name] = {
                                "headers": headers,
                                "rows": rows,
                                "numeric_cells": numeric_cells,
                            }
                    except Exception as exc:
                        st.write(f"    ❌ Erreur: {exc}")
            elif lot2_image_files and not api_key_is_set():
                st.write("  ⚠️ Images ignorées (clé API manquante)")

            # Extraction PDF contexte (premier slide)
            context_text = st.session_state.get("lot2_context_text", "")
            if lot2_context_pdf:
                st.write("  📄 Extraction du PDF de contexte...")
                try:
                    from parsers.pdf_questionnaire_parser import PDFQuestionnaireParser
                    ctx_bytes = lot2_context_pdf.read()
                    ctx_parser = PDFQuestionnaireParser()
                    context_text = ctx_parser.extract_text(ctx_bytes)
                    st.session_state["lot2_context_text"] = context_text
                    n_chars = len(context_text)
                    st.write(f"    ✅ {n_chars:,} caractères extraits du PDF contexte")
                except Exception as exc:
                    st.write(f"    ❌ Erreur extraction contexte : {exc}")

            # Parse questionnaire (PDF ou JSON)
            questionnaire_data = {}
            if lot2_questionnaire_file:
                ext = lot2_questionnaire_file.name.rsplit(".", 1)[-1].lower()

                if ext == "pdf":
                    st.write("  📋 Extraction du questionnaire PDF (Claude Vision)...")
                    try:
                        from parsers.pdf_questionnaire_parser import PDFQuestionnaireParser
                        pdf_bytes = lot2_questionnaire_file.read()
                        pdf_q_parser = PDFQuestionnaireParser(
                            api_key=get_api_key() if api_key_is_set() else None,
                            model=get_model(),
                        )
                        questionnaire_raw = pdf_q_parser.parse(
                            pdf_bytes,
                            use_llm=api_key_is_set(),
                        )
                        st.write(f"    🔍 {len(questionnaire_raw)} champ(s) détecté(s) dans le PDF")
                        # Affiche un aperçu
                        with st.expander("Aperçu questionnaire extrait du PDF", expanded=False):
                            st.json(questionnaire_raw)
                        q_parser = QuestionnaireParser()
                        questionnaire_data = q_parser.parse(questionnaire_raw)
                        total_q = sum(len(v) for v in questionnaire_data.values())
                        st.write(f"    ✅ {total_q} réponse(s) structurée(s)")
                    except Exception as exc:
                        parse_errors.append(f"Questionnaire PDF: {exc}")
                        st.write(f"    ❌ Erreur: {exc}")

                else:  # JSON
                    st.write("  📋 Parsing questionnaire JSON...")
                    try:
                        raw_json = lot2_questionnaire_file.read().decode("utf-8")
                        questionnaire_raw = json.loads(raw_json)
                        q_parser = QuestionnaireParser()
                        questionnaire_data = q_parser.parse(questionnaire_raw)
                        total_q = sum(len(v) for v in questionnaire_data.values())
                        st.write(f"    ✅ {total_q} réponse(s) parsées")
                    except Exception as exc:
                        parse_errors.append(f"Questionnaire JSON: {exc}")
                        st.write(f"    ❌ Erreur: {exc}")
            elif "questionnaire_manual_input" in st.session_state:
                manual_json_str = st.session_state.get("questionnaire_manual_input", "{}")
                if manual_json_str.strip() not in ("{}", ""):
                    try:
                        questionnaire_raw = json.loads(manual_json_str)
                        q_parser = QuestionnaireParser()
                        questionnaire_data = q_parser.parse(questionnaire_raw)
                        total_q = sum(len(v) for v in questionnaire_data.values())
                        st.write(f"  📋 Questionnaire manuel: {total_q} réponse(s)")
                    except Exception as exc:
                        st.write(f"  ⚠️ Questionnaire manuel invalide: {exc}")

            if parse_errors:
                status.update(
                    label=f"Étape 1 — Terminée avec {len(parse_errors)} erreur(s)",
                    state="error",
                )
            else:
                status.update(label="Étape 1 — Parsing terminé ✅", state="complete")

        # ===================================================================
        # STEP 2 — Compute KPIs
        # ===================================================================
        with st.status("Étape 2 — Calcul des KPIs...", expanded=True) as status:
            st.write("Initialisation du moteur KPI...")

            try:
                # Enrich raw_data with questionnaire numeric values
                if questionnaire_data.get("numerique"):
                    quant_sheet = {"headers": ["indicateur", "valeur"], "rows": [], "numeric_cells": []}
                    for i, (qkey, qval) in enumerate(questionnaire_data["numerique"].items(), start=2):
                        v = qval.get("valeur")
                        quant_sheet["rows"].append([qkey, v])
                        if isinstance(v, (int, float)):
                            quant_sheet["numeric_cells"].append({
                                "ref": f"B{i}",
                                "valeur": float(v),
                                "sheet": "Questionnaire",
                                "row": i,
                                "col": 2,
                            })
                    all_excel_data["sheets"]["Questionnaire"] = quant_sheet

                # Enrichir les règles KPI avec le mapping de la méthodologie
                methodology_for_mapping = st.session_state.get("methodology_text", "")
                methodo_mapping = parse_mapping_from_methodology(methodology_for_mapping)
                if methodo_mapping:
                    st.write(f"  📐 Mapping méthodologie détecté — {len(methodo_mapping)} KPI(s) enrichis")
                    from engine.kpi_engine import DEFAULT_RULES
                    enriched_rules = enrich_kpi_rules(DEFAULT_RULES, methodo_mapping)
                else:
                    enriched_rules = None
                    st.write("  ℹ️ Pas de mapping méthodologie — utilisation des hints par défaut")

                kpi_engine = KPIEngine(
                    raw_data=all_excel_data,
                    rules=enriched_rules,
                )
                kpis = kpi_engine.compute_all()
                st.session_state["lot2_kpis"] = kpis

                df_kpis = kpi_engine.get_as_dataframe()

                # Count by status
                statut_counts = df_kpis["statut"].value_counts()
                n_bon     = statut_counts.get("bon", 0)
                n_moyen   = statut_counts.get("moyen", 0)
                n_faible  = statut_counts.get("faible", 0)
                n_inconnu = statut_counts.get("inconnu", 0)

                st.write(
                    f"✅ {len(kpis)} KPIs calculés: "
                    f"🟢 {n_bon} bons | 🟡 {n_moyen} moyens | "
                    f"🔴 {n_faible} faibles | ⚪ {n_inconnu} inconnus"
                )

                # Display KPI table
                display_df = df_kpis.copy()
                display_df["statut"] = display_df["statut"].apply(format_kpi_statut)
                st.dataframe(
                    display_df[["label_fr", "valeur", "unite", "statut", "onglet", "cellule"]],
                    use_container_width=True,
                    hide_index=True,
                    column_config={
                        "label_fr": st.column_config.TextColumn("Indicateur"),
                        "valeur": st.column_config.NumberColumn("Valeur", format="%.2f"),
                        "unite": st.column_config.TextColumn("Unité"),
                        "statut": st.column_config.TextColumn("Statut"),
                        "onglet": st.column_config.TextColumn("Onglet source"),
                        "cellule": st.column_config.TextColumn("Cellule"),
                    },
                )

                # ── Checklist de disponibilité des données ────────────────
                st.write("  📋 Vérification de la couverture des données par slide...")
                data_check = check_slides_data(
                    kpi_dict=kpis,
                    methodology_mapping=methodo_mapping,
                    context_text=st.session_state.get("lot2_context_text", ""),
                )
                st.session_state["lot2_data_check"] = data_check

                n_ok_check = data_check["summary"].get(STATUS_OK, 0)
                n_ko_check = (
                    data_check["summary"].get(STATUS_MISSING, 0)
                    + data_check["summary"].get(STATUS_PARTIAL, 0)
                )
                if n_ko_check > 0:
                    st.write(
                        f"  ⚠️ {n_ko_check} slide(s) avec données insuffisantes — "
                        f"voir la checklist dans les résultats"
                    )
                else:
                    st.write(f"  ✅ Couverture données : {n_ok_check} slides complètes")

                # Auto-save KPIs dans le projet actif
                pid = st.session_state.get("current_project_id")
                if pid:
                    pm.save_kpis(pid, kpis)

                n_warn = data_check["summary"].get(STATUS_MISSING, 0) + data_check["summary"].get(STATUS_PARTIAL, 0)
                label_suffix = f" — {n_warn} slide(s) avec données insuffisantes ⚠️" if n_warn else " ✅"
                status.update(label=f"Étape 2 — KPIs calculés{label_suffix}", state="complete")

            except Exception as exc:
                st.error(f"Erreur lors du calcul des KPIs: {exc}")
                status.update(label=f"Étape 2 — Erreur: {exc}", state="error")
                st.stop()

        # ===================================================================
        # STEP 3 — Generate narrative
        # ===================================================================
        with st.status("Étape 3 — Génération du narratif...", expanded=True) as status:
            if not api_key_is_set():
                st.warning(
                    "⚠️ Clé API manquante — La génération LLM est ignorée. "
                    "Les diapositives seront créées avec des placeholders."
                )
                # Create placeholder slides
                from generation.llm_generator import PERFORMANCE_GLOBALE_SLIDES
                placeholder_slides = []
                for slide_def in PERFORMANCE_GLOBALE_SLIDES:
                    placeholder_slides.append({
                        "slide_id": slide_def["slide_id"],
                        "titre": slide_def["titre_defaut"],
                        "contenu": (
                            "[Contenu à générer — configurez la clé API Anthropic "
                            "pour activer la génération automatique]"
                        ),
                        "chiffres_cites": [],
                        "sources": [],
                        "erreur": None,
                    })
                st.session_state["lot2_slides"] = placeholder_slides
                status.update(
                    label="Étape 3 — Génération ignorée (pas de clé API)",
                    state="complete",
                )
            else:
                # Load methodology — priorité : session > bibliothèque active > fichier global
                methodology = st.session_state.get("methodology_text", "")
                if not methodology:
                    active_mid = st.session_state.get("methodo_active_id")
                    if active_mid:
                        methodology = ml.get_content(active_mid)
                    if not methodology:
                        methodology = load_methodology()
                if not methodology:
                    st.warning(
                        "⚠️ Méthodologie non définie — utilisation des instructions par défaut. "
                        "Allez dans 'Méthodologie' pour en sélectionner ou générer une."
                    )

                try:
                    generator = LLMGenerator(
                        api_key=get_api_key(),
                        model=get_model(),
                    )
                    kpis = st.session_state["lot2_kpis"]

                    st.write("Génération de 6 diapositives Performance Globale...")
                    slides = generator.generate_performance_globale(
                        kpi_dict=kpis,
                        methodology=methodology,
                        pharmacy_name=pharmacy_name,
                        context_text=st.session_state.get("lot2_context_text", ""),
                        image_results=st.session_state.get("lot1_image_results", []),
                    )
                    st.session_state["lot2_slides"] = slides

                    # Show slide summaries
                    for slide in slides:
                        if slide.get("erreur"):
                            st.write(f"  ❌ {slide['slide_id']}: {slide['erreur']}")
                        else:
                            n_chiffres = len(slide.get("chiffres_cites", []))
                            st.write(
                                f"  ✅ {slide['slide_id']}: "
                                f"'{slide['titre']}' — {n_chiffres} chiffre(s) cité(s)"
                            )

                    # Auto-save slides dans le projet actif
                    pid = st.session_state.get("current_project_id")
                    if pid:
                        pm.save_slides(pid, slides)
                        pm.update_project_meta(pid, statut="lot2")

                    status.update(
                        label=f"Étape 3 — {len(slides)} diapositives générées ✅",
                        state="complete",
                    )

                except Exception as exc:
                    st.error(f"Erreur lors de la génération: {exc}")
                    status.update(label=f"Étape 3 — Erreur: {exc}", state="error")
                    st.stop()

        # ===================================================================
        # STEP 4 — Audit
        # ===================================================================
        with st.status("Étape 4 — Audit anti-hallucination...", expanded=True) as status:
            slides = st.session_state.get("lot2_slides", [])
            kpis = st.session_state.get("lot2_kpis", {})

            if not slides:
                st.warning("Aucune diapositive à auditer.")
                status.update(label="Étape 4 — Ignorée", state="complete")
            else:
                audit_engine = AuditEngine()

                # Concatenate all generated content for audit
                all_content = "\n\n".join(
                    f"{s.get('titre', '')}\n{s.get('contenu', '')}"
                    for s in slides
                    if not s.get("erreur")
                )

                try:
                    audit_report = audit_engine.audit(
                    all_content,
                    kpis,
                    context_text=st.session_state.get("lot2_context_text", ""),
                    methodology_text=st.session_state.get("methodology_text", ""),
                )
                    st.session_state["lot2_audit"] = audit_report

                    # Display audit metrics
                    col1, col2, col3, col4 = st.columns(4)
                    with col1:
                        st.metric("Score audit", f"{audit_report['score_pct']:.1f}%")
                    with col2:
                        st.metric(
                            "Nombres validés",
                            f"{audit_report['validated']}/{audit_report['total_numbers_found']}",
                        )
                    with col3:
                        st.metric("Rejetés", len(audit_report["rejected"]))
                    with col4:
                        result_label = "✅ PASSÉ" if audit_report["passed"] else "❌ ÉCHOUÉ"
                        st.metric("Résultat", result_label)

                    if audit_report["passed"]:
                        st.success(f"✅ {audit_report['message']}")
                        status.update(
                            label=f"Étape 4 — Audit réussi ({audit_report['score_pct']:.1f}%) ✅",
                            state="complete",
                        )
                    else:
                        st.error(f"❌ {audit_report['message']}")
                        status.update(
                            label=f"Étape 4 — Audit échoué ({audit_report['score_pct']:.1f}%)",
                            state="error",
                        )

                except Exception as exc:
                    st.error(f"Erreur lors de l'audit: {exc}")
                    status.update(label=f"Étape 4 — Erreur: {exc}", state="error")

        # ===================================================================
        # STEP 5 — Build PPTX
        # ===================================================================
        with st.status("Étape 5 — Assemblage du PowerPoint...", expanded=True) as status:
            slides = st.session_state.get("lot2_slides", [])
            audit_report = st.session_state.get("lot2_audit", {})

            if not slides:
                st.warning("Aucune diapositive à assembler.")
                status.update(label="Étape 5 — Ignorée", state="complete")
            else:
                try:
                    assembler = PPTXAssembler()
                    pptx_bytes = assembler.build(
                        slides_content=slides,
                        pharmacy_name=pharmacy_name,
                    )
                    st.session_state["lot2_pptx_bytes"] = pptx_bytes

                    size_kb = len(pptx_bytes) / 1024
                    st.write(f"✅ PowerPoint assemblé ({size_kb:.1f} Ko, {len(slides)+1} diapositive(s))")

                    # Auto-save PPTX dans le projet actif
                    pid = st.session_state.get("current_project_id")
                    if pid:
                        pm.save_pptx(pid, pptx_bytes)
                        st.write("💾 PPTX sauvegardé dans le projet")

                    status.update(label="Étape 5 — PowerPoint prêt ✅", state="complete")

                except Exception as exc:
                    st.error(f"Erreur lors de l'assemblage PPTX: {exc}")
                    status.update(label=f"Étape 5 — Erreur: {exc}", state="error")

    # -----------------------------------------------------------------------
    # Results display (persistent, outside the pipeline button block)
    # -----------------------------------------------------------------------
    st.markdown("---")

    audit_report  = st.session_state.get("lot2_audit")
    slides        = st.session_state.get("lot2_slides")
    pptx_bytes    = st.session_state.get("lot2_pptx_bytes")
    data_check    = st.session_state.get("lot2_data_check")
    current_pharmacy = st.session_state.get("lot2_pharmacy_name", "pharmacie")

    # ── Checklist données (affiché dès que les KPIs sont calculés) ──────────
    if data_check:
        st.markdown(
            "<div style='font-size:0.65rem;font-weight:700;color:#5D6D7E;text-transform:uppercase;"
            "letter-spacing:0.12em;margin:0.5rem 0 0.75rem;'>📋 &nbsp;Checklist des données par slide</div>",
            unsafe_allow_html=True,
        )
        _render_data_checklist(data_check, expanded=(not bool(slides)))

    if audit_report and slides:
        st.markdown(
            "<div style='font-size:0.65rem;font-weight:700;color:#5D6D7E;text-transform:uppercase;"
            "letter-spacing:0.12em;margin:1.5rem 0 1rem;'>📊 &nbsp;Résultats du pipeline</div>",
            unsafe_allow_html=True,
        )

        # Audit summary
        if audit_report.get("passed", True):
            st.success(
                f"✅ Audit réussi — Score: {audit_report.get('score_pct', 100):.1f}%"
            )
        else:
            st.error(
                f"❌ Audit échoué — Score: {audit_report.get('score_pct', 0):.1f}% "
                f"— {len(audit_report.get('rejected', []))} nombre(s) halluciné(s) détecté(s)"
            )

            # Show rejected numbers
            if audit_report.get("rejected"):
                st.markdown("**Nombres non validés :**")
                for rejected in audit_report["rejected"]:
                    st.error(
                        f"**{rejected['number']}** — Contexte: `...{rejected['context']}...`"
                    )

        # Slide content preview
        with st.expander("👁️ Aperçu des diapositives générées"):
            for i, slide in enumerate(slides, start=1):
                st.markdown(f"### Diapositive {i}: {slide.get('titre', '(sans titre)')}")
                if slide.get("erreur"):
                    st.error(f"Erreur: {slide['erreur']}")
                else:
                    st.markdown(slide.get("contenu", ""))
                    if slide.get("chiffres_cites"):
                        st.caption(
                            f"Chiffres cités: {', '.join(str(n) for n in slide['chiffres_cites'])}"
                        )
                st.markdown("---")

        # Download PPTX
        if pptx_bytes:
            safe_name = "".join(
                c if c.isalnum() or c in "_ -" else "_"
                for c in current_pharmacy
            ).strip()
            filename = f"AUDIT_360_{safe_name or 'Pharmacie'}_STRATEGIE_Part2.pptx"

            st.download_button(
                label="📥 Télécharger le PowerPoint (Performance Globale)",
                data=pptx_bytes,
                file_name=filename,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary",
                use_container_width=True,
            )

        # KPI dataframe
        kpis = st.session_state.get("lot2_kpis")
        if kpis:
            with st.expander("📈 Détail des KPIs calculés"):
                kpi_engine = KPIEngine(raw_data={"sheets": {}, "source": ""})
                kpi_engine._kpis = kpis
                df = kpi_engine.get_as_dataframe()
                df["statut"] = df["statut"].apply(format_kpi_statut)
                st.dataframe(
                    df,
                    use_container_width=True,
                    hide_index=True,
                )
